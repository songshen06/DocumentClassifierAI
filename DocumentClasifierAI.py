import os
import glob
import warnings
import argparse
import csv
from openai import OpenAI
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadWarning
from pdfminer.high_level import extract_text as pdfminer_extract_text
from pptx import Presentation
from docx import Document
from docx import Document as DocxDocument
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Define the path where documents are downloaded
download_path = os.path.expanduser("/Users/shens/Downloads")

# Model definitions for easy modification
OLLAMA_MODEL = 'gemma2:latest'
DEEPSEEK_MODEL = 'deepseek-chat'
NVIDIA_MODEL = 'mistralai/mistral-large-2-instruct'

# Ignore PdfReadWarning
warnings.filterwarnings("ignore", category=PdfReadWarning)

class OpenAIWrapper:
    def __init__(self, provider="ollama"):
        if provider == "ollama":
            self.client = OpenAI(
                base_url='http://localhost:11434/v1',
                api_key='ollama'
            )
            self.model = OLLAMA_MODEL
        elif provider == "deepseek":
            self.client = OpenAI(
                api_key=os.getenv("DEEPKEY"),
                base_url="https://api.deepseek.com"
            )
            self.model = DEEPSEEK_MODEL
        elif provider == "nvidia_nim":
            self.client = OpenAI(
                base_url="https://integrate.api.nvidia.com/v1",
                api_key=os.getenv("NVIDIA_key")
            )
            self.model = NVIDIA_MODEL
        else:
            raise ValueError(f"Unknown provider: {provider}")

    def create_completion(self, messages, stream=False):
        return self.client.chat.completions.create(
            model=self.model,
            messages=messages,
            stream=stream
        )

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        text = pdfminer_extract_text(pdf_path)
        text = text[:1000]  # Limit the text to the first 1000 characters to avoid exceeding token limits
    except Exception as e:
        print(f"Error reading PDF {pdf_path}: {e}")
    return text

def extract_text_from_docx(doc_path):
    text = ""
    try:
        document = Document(doc_path)
        for para in document.paragraphs:
            text += para.text + "\n"
        text = text[:1000]  # Limit the text to the first 1000 characters to avoid exceeding token limits
    except Exception as e:
        print(f"Error reading DOCX {doc_path}: {e}")
    return text

def extract_text_from_pptx(ppt_path):
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        text = text[:1000]  # Limit the text to the first 1000 characters to avoid exceeding token limits
    except Exception as e:
        print(f"Error reading PPTX {ppt_path}: {e}")
    return text

def summarize_document(text):
    # Use ollama model to summarize the document
    ollama_wrapper = OpenAIWrapper(provider="ollama")
    summary_response = ollama_wrapper.create_completion(
        messages=[
            {"role": "system", "content": "You are an assistant that summarizes documents. Provide a concise summary of the following document content:"},
            {"role": "user", "content": text}
        ],
        stream=False
    )
    return summary_response.choices[0].message.content.strip()

def classify_documents(product_keywords, document_paths, provider="ollama"):
    ai_wrapper = OpenAIWrapper(provider=provider)
    classified_documents = []

    total_documents = len(document_paths)
    print(f"Total documents to process: {total_documents}")

    for index, doc_path in enumerate(document_paths, start=1):
        print(f"Processing file {index}/{total_documents}: {doc_path}")
        if doc_path.endswith(".pdf"):
            text = extract_text_from_pdf(doc_path)
        elif doc_path.endswith(".docx"):
            text = extract_text_from_docx(doc_path)
        elif doc_path.endswith(".pptx"):
            text = extract_text_from_pptx(doc_path)
        else:
            continue

        # Generate summary of the document using ollama
        summary = summarize_document(text)

        # Use the summary for classification with the chosen provider
        response = ai_wrapper.create_completion(
            messages=[
                {"role": "system", "content": (
                    "You are an assistant that classifies documents into the following categories: {}. "
                    "Please only return the category name without any additional information. The categories are: Omniverse, vGPU, NVAIE, Uncategorized.\n"
                    "For example, if the document discusses 3D rendering and digital twins, return 'Omniverse'. If the document discusses GPU virtualization, remote workstation, or virtualized graphics, return 'vGPU'.\n"
                    "Make sure to only return one of the following words: Omniverse, vGPU, NVAIE, Uncategorized. No other text should be included in the response."
                ).format(list(product_keywords.keys()))},
                {"role": "user", "content": summary}
            ],
            stream=False
        )

        category = response.choices[0].message.content.strip()
        classified_documents.append((doc_path, category))

        remaining_documents = total_documents - index
        print(f"Remaining documents: {remaining_documents}")
        print(f"Document '{doc_path}' classified as: {category}")

    return classified_documents

def save_classification_to_csv(classified_documents, output_path, output_filename):
    csv_path = os.path.join(output_path, output_filename)
    with open(csv_path, mode='w', newline='') as file:
        writer = csv.writer(file)
        writer.writerow(["File Name", "Category"])
        for doc_path, category in classified_documents:
            writer.writerow([doc_path, category])
    print(f"Classification results saved to {csv_path}")

def main():
    parser = argparse.ArgumentParser(description="Classify documents based on their content.")
    parser.add_argument("-FT", "--file_type", type=str, choices=["pdf", "docx", "pptx", "all"], default="all", help="Type of files to process: pdf, docx, pptx, or all")
    parser.add_argument("-OD", "--output_dir", type=str, default=".", help="Directory to save the output file. Defaults to the current directory.")
    parser.add_argument("-OF", "--output_filename", type=str, default="classification_results.csv", help="Name of the output CSV file. Defaults to 'classification_results.csv'.")
    parser.add_argument("-P", "--provider", type=str, choices=["ollama", "deepseek", "nvidia_nim"], default="ollama", help="LLM provider to use for classification: ollama, deepseek, or nvidia_nim")
    args = parser.parse_args()

    # Input validation for output directory
    if not os.path.exists(args.output_dir):
        print(f"Error: The specified output directory '{args.output_dir}' does not exist.")
        return

    # Input validation for output filename
    if not args.output_filename.endswith(".csv"):
        print(f"Error: The output filename '{args.output_filename}' must have a .csv extension.")
        return

    # Define the product keywords/categories
    product_keywords = {
        "Omniverse": ["omniverse", "Omniverse", "digital twin", "rendering", "Kit", "Nucleus"],
        "vGPU": [
            "vGPU", "vWS", "vPC", "virtual GPU", "GPU sharing", "remote workstation", 
            "vGPU licensing", "NVIDIA GRID", "NVIDIA vComputeServer", "virtualized graphics",
            "GPU acceleration", "remote desktop", "vPC profiles", "vApps", "virtual workstation",
            "graphics virtualization"
        ],
        "NVAIE": ["Nvidia enterprise", "NVAIE", "AI infrastructure", "cloud services", "NVIDIA AI"]
    }

    document_paths = []
    if args.file_type in ["pdf", "all"]:
        document_paths += [doc_path for doc_path in glob.glob(os.path.join(download_path, "**", "*.pdf"), recursive=True) if not os.path.basename(doc_path).startswith('$')]
    if args.file_type in ["docx", "all"]:
        document_paths += [doc_path for doc_path in glob.glob(os.path.join(download_path, "**", "*.docx"), recursive=True) if not os.path.basename(doc_path).startswith('$')]
    if args.file_type in ["pptx", "all"]:
        document_paths += [doc_path for doc_path in glob.glob(os.path.join(download_path, "**", "*.pptx"), recursive=True) if not os.path.basename(doc_path).startswith('$')]

    # Check if any documents were found
    if not document_paths:
        print("Error: No documents found for the specified file type.")
        return

    # Print the number of files by type
    pdf_count = len(glob.glob(os.path.join(download_path, "**", "*.pdf"), recursive=True)) if args.file_type in ["pdf", "all"] else 0
    docx_count = len(glob.glob(os.path.join(download_path, "**", "*.docx"), recursive=True)) if args.file_type in ["docx", "all"] else 0
    pptx_count = len(glob.glob(os.path.join(download_path, "**", "*.pptx"), recursive=True)) if args.file_type in ["pptx", "all"] else 0
    print(f"PDF files: {pdf_count}, DOCX files: {docx_count}, PPTX files: {pptx_count}")

    classified_documents = classify_documents(product_keywords, document_paths, provider=args.provider)

    # Save the classification results to a CSV file
    save_classification_to_csv(classified_documents, args.output_dir, args.output_filename)

if __name__ == "__main__":
    main()
